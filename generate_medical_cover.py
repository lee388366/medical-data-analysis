import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 用户可自定义内容
HOSPITAL = "中国医学科学院整形外科医院"
DEPARTMENT = "急危重症科"
APPLICANT = "张三｜主治医师"
AFFILIATION = "中国医学科学院整形外科医院"
DATE = "2025年6月"
LOGO_PATH = "logo.png"  # 如有医院logo图片，放在同目录并改名

# 颜色方案
BG_COLOR = RGBColor(245, 248, 250)      # #F5F8FA
TITLE_COLOR = RGBColor(0, 78, 138)      # #004E8A
SUB_COLOR = RGBColor(0, 78, 138)        # #004E8A
TEXT_COLOR = RGBColor(51, 51, 51)       # #333333
ACCENT_COLOR = RGBColor(0, 160, 233)    # #00A0E9

prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页

# 背景色
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = BG_COLOR

# Logo（如有）
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(0.5), width=Inches(1.5))

# 主标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), slide_width - Inches(1), Inches(1.5))
tf = title_box.text_frame
tf.clear()
p = tf.add_paragraph()
p.text = HOSPITAL
p.font.bold = True
p.font.size = Pt(56)
p.font.name = 'Microsoft YaHei'
p.font.color.rgb = TITLE_COLOR
p.alignment = PP_ALIGN.CENTER

# 副标题
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), slide_width - Inches(1), Inches(1.0))
tf2 = subtitle_box.text_frame
tf2.clear()
p2 = tf2.add_paragraph()
p2.text = DEPARTMENT
p2.font.bold = True
p2.font.size = Pt(40)
p2.font.name = 'Microsoft YaHei'
p2.font.color.rgb = SUB_COLOR
p2.alignment = PP_ALIGN.CENTER

# 强调横线
line_left = (slide_width - Inches(2)) // 2
line_top = Inches(4.1)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, Inches(2), Pt(8))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_COLOR
line.line.width = Pt(0)

# 信息区
info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), slide_width - Inches(1), Inches(1.5))
tf3 = info_box.text_frame
tf3.clear()
p3 = tf3.add_paragraph()
p3.text = APPLICANT
p3.font.bold = True
p3.font.size = Pt(32)
p3.font.name = 'Microsoft YaHei'
p3.font.color.rgb = TEXT_COLOR
p3.alignment = PP_ALIGN.CENTER

p4 = tf3.add_paragraph()
p4.text = AFFILIATION
p4.font.size = Pt(28)
p4.font.name = 'Microsoft YaHei'
p4.font.color.rgb = TEXT_COLOR
p4.alignment = PP_ALIGN.CENTER

p5 = tf3.add_paragraph()
p5.text = DATE
p5.font.size = Pt(28)
p5.font.name = 'Microsoft YaHei'
p5.font.color.rgb = TEXT_COLOR
p5.alignment = PP_ALIGN.CENTER

prs.save('medical_cover.pptx')
print('已生成 medical_cover.pptx，可用WPS/PowerPoint打开。') 